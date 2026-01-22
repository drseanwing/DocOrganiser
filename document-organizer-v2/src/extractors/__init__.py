"""
Text Extractors - Extract text content from various file types.

Supports:
- Plain text (txt, md, csv, html, json)
- PDF (using PyMuPDF)
- Word documents (using python-docx)
- Excel spreadsheets (using openpyxl)
- PowerPoint presentations (using python-pptx)
"""

import asyncio
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional
import structlog

logger = structlog.get_logger("extractors")


class BaseExtractor(ABC):
    """Abstract base class for text extractors."""
    
    @abstractmethod
    async def extract(self, file_path: Path) -> Optional[str]:
        """
        Extract text content from a file.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text or None on failure
        """
        pass


class TextExtractor(BaseExtractor):
    """Extractor for plain text files (txt, md, csv, html, json, xml)."""
    
    async def extract(self, file_path: Path) -> Optional[str]:
        try:
            # Try common encodings
            for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # Fall back to binary read with chardet
            import chardet
            with open(file_path, 'rb') as f:
                raw = f.read()
                detected = chardet.detect(raw)
                encoding = detected.get('encoding', 'utf-8')
                return raw.decode(encoding, errors='replace')
                
        except Exception as e:
            logger.warning("text_extraction_failed", path=str(file_path), error=str(e))
            return None


class PDFExtractor(BaseExtractor):
    """Extractor for PDF files using PyMuPDF (fitz)."""
    
    async def extract(self, file_path: Path) -> Optional[str]:
        try:
            import fitz  # PyMuPDF
            
            text_parts = []
            with fitz.open(str(file_path)) as doc:
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{text}")
            
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.warning("pymupdf_not_installed")
            return await self._fallback_pdftotext(file_path)
        except Exception as e:
            logger.warning("pdf_extraction_failed", path=str(file_path), error=str(e))
            return None
    
    async def _fallback_pdftotext(self, file_path: Path) -> Optional[str]:
        """Fallback to pdftotext command-line tool."""
        try:
            proc = await asyncio.create_subprocess_exec(
                'pdftotext', '-layout', str(file_path), '-',
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await proc.communicate()
            if proc.returncode == 0:
                return stdout.decode('utf-8', errors='replace')
        except Exception as e:
            logger.warning("pdftotext_failed", error=str(e))
        return None


class DocxExtractor(BaseExtractor):
    """Extractor for Word documents using python-docx."""
    
    async def extract(self, file_path: Path) -> Optional[str]:
        try:
            from docx import Document
            
            doc = Document(str(file_path))
            
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n".join(table_text))
            
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.warning("python_docx_not_installed")
            return await self._fallback_pandoc(file_path)
        except Exception as e:
            logger.warning("docx_extraction_failed", path=str(file_path), error=str(e))
            return None
    
    async def _fallback_pandoc(self, file_path: Path) -> Optional[str]:
        """Fallback to pandoc for conversion."""
        try:
            proc = await asyncio.create_subprocess_exec(
                'pandoc', '-f', 'docx', '-t', 'plain', str(file_path),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await proc.communicate()
            if proc.returncode == 0:
                return stdout.decode('utf-8', errors='replace')
        except Exception as e:
            logger.warning("pandoc_failed", error=str(e))
        return None


class XlsxExtractor(BaseExtractor):
    """Extractor for Excel spreadsheets using openpyxl."""
    
    async def extract(self, file_path: Path) -> Optional[str]:
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(str(file_path), data_only=True, read_only=True)
            
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]
                
                row_count = 0
                for row in sheet.iter_rows(max_row=100, values_only=True):  # Limit rows
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        sheet_text.append(" | ".join(row_values))
                        row_count += 1
                
                if row_count > 0:
                    text_parts.append("\n".join(sheet_text))
            
            wb.close()
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.warning("openpyxl_not_installed")
            return None
        except Exception as e:
            logger.warning("xlsx_extraction_failed", path=str(file_path), error=str(e))
            return None


class PptxExtractor(BaseExtractor):
    """Extractor for PowerPoint presentations using python-pptx."""
    
    async def extract(self, file_path: Path) -> Optional[str]:
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.warning("python_pptx_not_installed")
            return None
        except Exception as e:
            logger.warning("pptx_extraction_failed", path=str(file_path), error=str(e))
            return None


# Extractor registry
EXTRACTORS = {
    # Plain text
    'txt': TextExtractor(),
    'md': TextExtractor(),
    'csv': TextExtractor(),
    'html': TextExtractor(),
    'htm': TextExtractor(),
    'json': TextExtractor(),
    'xml': TextExtractor(),
    'log': TextExtractor(),
    'ini': TextExtractor(),
    'cfg': TextExtractor(),
    'yaml': TextExtractor(),
    'yml': TextExtractor(),
    
    # Documents
    'pdf': PDFExtractor(),
    'docx': DocxExtractor(),
    'doc': DocxExtractor(),  # Will use pandoc fallback
    
    # Spreadsheets
    'xlsx': XlsxExtractor(),
    'xls': XlsxExtractor(),  # May not work perfectly
    
    # Presentations
    'pptx': PptxExtractor(),
    'ppt': PptxExtractor(),  # May not work
}


def get_extractor(extension: str) -> Optional[BaseExtractor]:
    """
    Get the appropriate extractor for a file extension.
    
    Args:
        extension: File extension (without dot)
        
    Returns:
        Extractor instance or None if unsupported
    """
    return EXTRACTORS.get(extension.lower())


def is_supported(extension: str) -> bool:
    """Check if a file extension is supported for text extraction."""
    return extension.lower() in EXTRACTORS
